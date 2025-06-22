# app.py — Final with Upload, Chat, File Manager, Memory

import streamlit as st
import os
import json
import tempfile
from PIL import Image
import pytesseract
from pdf2image import convert_from_path
import fitz  # PyMuPDF
import docx
import pptx
import pandas as pd
from datetime import datetime
from langchain.text_splitter import CharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain.embeddings import HuggingFaceEmbeddings
from transformers import pipeline

# ---- CONFIG ----
st.set_page_config(page_title="🧠 SmartDoc Chatbot", layout="wide")
UPLOAD_DIR = "uploads"
MEMORY_DIR = "memory"
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(MEMORY_DIR, exist_ok=True)

# ---- LOAD LOGIN ----
def authenticate_user():
    with open("auth.json") as f:
        auth_data = json.load(f)
    login_container = st.empty()
    with login_container.form("Login"):
        st.subheader("🔐 Login")
        user = st.text_input("User ID")
        pwd = st.text_input("Password", type="password")
        login = st.form_submit_button("Login")
        if login:
            if user in auth_data and auth_data[user] == pwd:
                st.session_state["user"] = user
                st.session_state["authenticated"] = True
                login_container.empty()
                return True
            else:
                st.error("❌ Invalid credentials")
    return False

# ---- ADVANCED FILE PARSER ----
def smart_parse(file, save=True):
    ext = file.name.split(".")[-1].lower()
    text = ""
    saved_path = os.path.join(UPLOAD_DIR, file.name)

    if save:
        with open(saved_path, "wb") as f:
            f.write(file.getbuffer())

    with tempfile.NamedTemporaryFile(delete=False) as tmp:
        tmp.write(file.read())
        tmp_path = tmp.name

    try:
        if ext == "pdf":
            doc = fitz.open(saved_path)
            for page in doc:
                text += page.get_text()
            if not text.strip():
                images = convert_from_path(saved_path)
                for img in images:
                    text += pytesseract.image_to_string(img)
        elif ext == "docx":
            doc = docx.Document(saved_path)
            text = "\n".join(p.text for p in doc.paragraphs)
        elif ext == "pptx":
            prs = pptx.Presentation(saved_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif ext == "xlsx":
            df = pd.read_excel(saved_path)
            text = df.astype(str).apply(" ".join, axis=1).str.cat(sep="\n")
        elif ext in ["jpg", "jpeg", "png"]:
            img = Image.open(saved_path)
            text = pytesseract.image_to_string(img)
        else:
            with open(saved_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
    except Exception as e:
        st.error(f"❌ Error parsing file: {e}")
    return text

# ---- VECTOR DB ----
def build_vectorstore(text):
    splitter = CharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
    docs = splitter.create_documents([text])
    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    return FAISS.from_documents(docs, embeddings)

# ---- INIT FLAN BOT ----
def setup_bot():
    return pipeline("text2text-generation", model="google/flan-t5-base")

# ---- CONVERSATION UI ----
def chat_interface():
    st.subheader("💬 Ask Your Questions")
    if "chat_history" not in st.session_state:
        st.session_state["chat_history"] = []

    query = st.text_input("Your question:")
    if query:
        vectorstore = st.session_state["vectorstore"]
        retriever = vectorstore.as_retriever()
        context_docs = retriever.get_relevant_documents(query)
        context = "\n".join([doc.page_content for doc in context_docs])

        history = "\n".join([f"Q: {q}\nA: {a}" for q, a in st.session_state["chat_history"][-3:]])

        prompt = f"""
You are a helpful assistant. Answer based on context and previous Q&A.

{history}

Context:
{context}

Question: {query}
Answer:
"""
        bot = st.session_state["bot"]
        answer = bot(prompt, max_length=512, do_sample=False)[0]['generated_text']
        st.session_state["chat_history"].append((query, answer))

        # Save chat to memory folder
        user = st.session_state.get("user", "guest")
        now = datetime.now().strftime("%Y-%m-%d")
        memory_file = os.path.join(MEMORY_DIR, f"{user}_chat_{now}.txt")
        with open(memory_file, "a", encoding="utf-8") as f:
            f.write(f"Q: {query}\nA: {answer}\n\n")

    for q, a in reversed(st.session_state["chat_history"]):
        st.markdown(f"**You:** {q}")
        st.markdown(f"**Bot:** {a}")

# ---- APP LOGIC ----
if "authenticated" not in st.session_state:
    st.session_state["authenticated"] = False

menu = st.sidebar.radio("📂 Menu", ["Chat", "Upload"])

if menu == "Upload":
    if not st.session_state["authenticated"]:
        if not authenticate_user():
            st.stop()

    st.subheader("📂 File Management")

    col1, col2 = st.columns(2)

    with col1:
        if st.button("📤 Upload New File"):
            st.session_state["show_upload"] = True
            st.session_state["show_manage"] = False

    with col2:
        if st.button("🗂️ Manage Uploaded Files"):
            st.session_state["show_upload"] = False
            st.session_state["show_manage"] = True

    # ---- Upload Panel ----
    if st.session_state.get("show_upload"):
        st.markdown("### 📤 Upload Documents")
        files = st.file_uploader("Upload any document", accept_multiple_files=True)
        if files:
            with st.spinner("Processing files..."):
                full_text = ""
                for file in files:
                    full_text += smart_parse(file, save=True) + "\n"
                st.session_state["vectorstore"] = build_vectorstore(full_text)
                st.session_state["bot"] = setup_bot()
            st.success("✅ Files processed and chatbot ready!")

    # ---- File Management Panel ----
    if st.session_state.get("show_manage"):
        st.markdown("### 🗂️ Uploaded Files")
        file_list = os.listdir(UPLOAD_DIR)
        if not file_list:
            st.info("No uploaded files found.")
        else:
            for filename in file_list:
                filepath = os.path.join(UPLOAD_DIR, filename)
                size = os.path.getsize(filepath) / 1024
                file_col1, file_col2 = st.columns([4, 1])
                with file_col1:
                    st.write(f"📄 `{filename}` — {size:.2f} KB")
                with file_col2:
                    if st.button("🗑️ Delete", key=filename):
                        os.remove(filepath)
                        st.success(f"Deleted `{filename}`")
                        st.experimental_rerun()

    # ---- Navigation Buttons ----
    st.divider()
    if st.button("⬅️ Back to Chat"):
        st.switch_page("/app.py")
    if st.button("🔒 Logout"):
        st.session_state.clear()
        st.experimental_rerun()

elif menu == "Chat":
    if "vectorstore" not in st.session_state or "bot" not in st.session_state:
        st.info("📥 Please upload files first from the Upload tab.")
        st.stop()
    chat_interface()


